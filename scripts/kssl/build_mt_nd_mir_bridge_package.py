from pathlib import Path
import pandas as pd
from docx import Document
from docx.shared import Inches as DI, Pt as DP
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

ROOT=Path(__file__).resolve().parents[2]
T=ROOT/'outputs/tables/kssl_mir_bridge'; F=ROOT/'outputs/figures/kssl_mir_bridge'
Q=ROOT/'outputs/figures/kssl_mir_qc'; A=ROOT/'outputs/figures/kssl_mir_analysis'
PPT=ROOT/'outputs/presentations/KSSL_MT_ND_MIR_bridge_results.pptx'
DOC=ROOT/'outputs/reports/KSSL_MT_ND_MIR_bridge_results.docx'
NAVY=RGBColor(23,50,77); TEAL=RGBColor(15,124,128); GOLD=RGBColor(217,164,65); GRAY=RGBColor(102,116,125)

def title(s,text,sub=''):
    b=s.shapes.add_textbox(Inches(.55),Inches(.32),Inches(12.2),Inches(.52)); r=b.text_frame.paragraphs[0].add_run(); r.text=text; r.font.size=Pt(24); r.font.bold=True; r.font.color.rgb=NAVY
    b=s.shapes.add_textbox(Inches(.57),Inches(.9),Inches(12),Inches(.3)); r=b.text_frame.paragraphs[0].add_run(); r.text=sub; r.font.size=Pt(11); r.font.color.rgb=GRAY
    z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.14)); z.fill.solid(); z.fill.fore_color.rgb=TEAL; z.line.fill.background()
def bullets(s,items,x,y,w,h,size=15):
    b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear()
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.size=Pt(size); p.font.color.rgb=NAVY; p.space_after=Pt(10)

def slides(metrics):
    m=metrics.set_index('target'); prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'A quality-controlled MT–ND MIR cohort is analysis-ready','398 surface samples; 1,592 selected scans; four technical replicates averaged per master')
    s.shapes.add_picture(str(Q/'mean_mir_spectra_by_spatial_group.png'),Inches(.55),Inches(1.3),width=Inches(7.5))
    bullets(s,['398 of 404 mapped samples have physical spectra','Common grid: 4000–600 cm⁻¹ at 2 cm⁻¹ spacing','11 rescanned samples resolved without double-counting','Six missing samples are documented, not silently dropped','Spatial evidence remains an external weak reference'],8.35,1.55,4.25,4.8,15)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'MIR variation is continuous—not a simple hydric split','SNV + first derivative PCA; exploratory and unsupervised')
    s.shapes.add_picture(str(A/'mir_pca_by_spatial_group.png'),Inches(.55),Inches(1.25),width=Inches(7.6))
    bullets(s,['PC1 explains 31.2%; PC2 explains 20.5%','Groups overlap substantially in spectral space','Overlap is expected because wetland status is not one soil constituent','A multivariate model is needed, with independent validation'],8.45,1.7,4.15,4.3,16)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'MIR robustly predicts key laboratory bridge properties','Five-fold cross-validation holds out entire KSSL projects')
    s.shapes.add_picture(str(F/'mir_grouped_cv_summary.png'),Inches(.55),Inches(1.35),width=Inches(7.5))
    bullets(s,[f"Total C: R²={m.loc['total_carbon_pct','r2_grouped_cv']:.2f}, ρ={m.loc['total_carbon_pct','spearman_rho']:.2f}",f"Clay: R²={m.loc['clay_pct','r2_grouped_cv']:.2f}, ρ={m.loc['clay_pct','spearman_rho']:.2f}",f"CEC: R²={m.loc['cec_nh4oac_cmol_kg','r2_grouped_cv']:.2f}, ρ={m.loc['cec_nh4oac_cmol_kg','spearman_rho']:.2f}",f"15-bar water: R²={m.loc['water_retention_15bar_pct','r2_grouped_cv']:.2f}, ρ={m.loc['water_retention_15bar_pct','spearman_rho']:.2f}",'MIR provides credible laboratory variables for later airborne linkage.'],8.3,1.45,4.4,5.1,15)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Spatial hydric evidence is detectable—but not deterministic','Ordinal evidence score: neither=0, one mapped source=1, both sources=2')
    s.shapes.add_picture(str(F/'mir_grouped_cv_observed_predicted.png'),Inches(.45),Inches(1.25),width=Inches(8.7))
    bullets(s,[f"Out-of-project ρ={m.loc['spatial_evidence_score','spearman_rho']:.2f}",f"Out-of-project R²={m.loc['spatial_evidence_score','r2_grouped_cv']:.2f}",'Signal persists beyond individual KSSL projects.','Performance is modest: this is supporting evidence, not a hydric classifier.','Next validation should hold out geography and use field-confirmed labels.'],9.35,1.55,3.45,4.9,14)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'How this advances the airborne-sensing goal','Laboratory MIR supplies interpretable soil-property targets and a tested evidence gradient')
    bullets(s,['Use MIR-derived carbon, clay, pH, water retention, CEC, and Fe as laboratory bridge variables.','Compare those variables with airborne VNIR–SWIR observations at exposed or minimally vegetated locations.','Control vegetation, moisture, spectral mixing, observation date, and surface-versus-subsurface mismatch.','Keep SSURGO/NWI evidence independent from chemistry and spectra.','Validate by project and geography before any delineation claim.'],.9,1.45,11.6,5.2,18)
    PPT.parent.mkdir(parents=True,exist_ok=True); prs.save(PPT)

def report(metrics):
    d=Document(); d.styles['Normal'].font.name='Aptos'; d.styles['Normal'].font.size=DP(10)
    d.add_heading('Montana–North Dakota MIR Bridge Analysis',0)
    d.add_paragraph('Objective: test whether quality-controlled KSSL MIR spectra recover laboratory soil properties and track independent spatial hydric evidence while preserving project-level independence.')
    d.add_heading('Cohort and preprocessing',1)
    d.add_paragraph('The analysis includes 398 of 404 mapped surface samples. Four technical replicates were interpolated to a common 4000–600 cm⁻¹ grid and averaged within the selected MIR master. Eleven rescanned samples were resolved by prioritizing passed scans, complete replicate sets, and the latest master. Six samples lacking physical files are documented separately.')
    d.add_picture(str(Q/'mean_mir_spectra_by_spatial_group.png'),width=DI(6.8))
    d.add_heading('Validation design',1)
    d.add_paragraph('PLS models used SNV-normalized first-derivative spectra and ten components. Five-fold cross-validation held out complete KSSL laboratory projects, reducing leakage from related samples and project-specific conditions. Component tuning was not optimized; results are an initial fixed-model benchmark.')
    table=d.add_table(rows=1,cols=6); table.style='Light Shading Accent 1'
    for c,x in zip(table.rows[0].cells,['Target','n','Projects','R²','RMSE','Spearman ρ']): c.text=x
    for _,r in metrics.iterrows():
        q=table.add_row().cells
        for c,x in zip(q,[r.label,f'{r.n:,}',str(r.project_groups),f'{r.r2_grouped_cv:.2f}',f'{r.rmse_grouped_cv:.2f}',f'{r.spearman_rho:.2f}']): c.text=x
    d.add_heading('Interpretation',1)
    d.add_paragraph('MIR predicts total carbon, clay, pH, 15-bar water retention, and CEC strongly outside the projects used for fitting. Iron measures are recoverable more modestly, partly reflecting smaller sample sizes. The ordinal spatial-evidence score is detectable but substantially less predictable, consistent with hydric status depending on landscape position, hydrology, morphology, and map scale rather than chemistry alone.')
    d.add_picture(str(F/'mir_grouped_cv_summary.png'),width=DI(6.8))
    d.add_heading('Boundaries on inference',1)
    for x in ['The evidence score is not a confirmed hydric/non-hydric label.','NWI and SSURGO differ in mapping purpose, resolution, and vintage.','The score assigns equal ordinal spacing to categories for an exploratory benchmark.','Project-grouped validation does not replace geographic or independent field validation.','Laboratory MIR spectra of dried, ground samples are not directly equivalent to airborne VNIR–SWIR reflectance.']: d.add_paragraph(x,style='List Bullet')
    d.add_heading('Next step',1); d.add_paragraph('Use the successfully predicted laboratory properties as bridge variables for airborne analysis, while pursuing field-confirmed hydric indicators and geographically independent validation sites. The immediate modeling extension should test geographic holdouts and uncertainty, not optimize a binary classifier against weak labels.')
    DOC.parent.mkdir(parents=True,exist_ok=True); d.save(DOC)

if __name__=='__main__':
    metrics=pd.read_csv(T/'grouped_cv_metrics.csv'); slides(metrics); report(metrics); print(PPT); print(DOC)

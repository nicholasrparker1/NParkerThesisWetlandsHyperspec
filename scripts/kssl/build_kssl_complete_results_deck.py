from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

R=Path(__file__).resolve().parents[2]
OUT=R/'outputs/presentations/KSSL_hydric_project_complete_results.pptx'
NAVY=RGBColor(23,50,77); TEAL=RGBColor(15,124,128); GOLD=RGBColor(217,164,65); GRAY=RGBColor(102,116,125); PALE=RGBColor(234,242,243)
def text(s,x,y,w,h,value,size=14,color=NAVY,bold=False,align=PP_ALIGN.LEFT):
 b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear(); tf.margin_left=tf.margin_right=0
 p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=value; r.font.name='Aptos'; r.font.size=Pt(size); r.font.bold=bold; r.font.color.rgb=color; return b
def title(s,t,sub=''):
 text(s,.55,.31,12.2,.55,t,24,NAVY,True); text(s,.57,.9,12,.3,sub,11,GRAY)
 z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.14)); z.fill.solid(); z.fill.fore_color.rgb=TEAL; z.line.fill.background()
def bullets(s,items,x,y,w,h,size=15):
 b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear()
 for i,item in enumerate(items):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.name='Aptos'; p.font.size=Pt(size); p.font.color.rgb=NAVY; p.space_after=Pt(9)
def image(s,p,x,y,w,h=None): s.shapes.add_picture(str(p),Inches(x),Inches(y),width=Inches(w),height=Inches(h) if h else None)
def footer(s,msg): text(s,.65,7.08,12,.22,msg,9,GRAY,False,PP_ALIGN.CENTER)
def section(s,t,sub):
 z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(7.5)); z.fill.solid(); z.fill.fore_color.rgb=NAVY; z.line.fill.background()
 text(s,.8,2.45,11.7,.7,t,30,RGBColor(255,255,255),True,PP_ALIGN.CENTER); text(s,1.2,3.35,10.9,.55,sub,17,RGBColor(210,225,228),False,PP_ALIGN.CENTER)
def slide(prs): return prs.slides.add_slide(prs.slide_layouts[6])
def main():
 prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
 s=slide(prs); title(s,'Integrating Laboratory Soil Spectroscopy with Spatial Wetland Evidence','Montana–North Dakota KSSL analysis toward airborne hydric-soil delineation')
 bullets(s,['Goal: determine how laboratory MIR spectra and measured soil properties can improve interpretation of remotely sensed wetland soils.','Current contribution: a traceable laboratory–spatial bridge built before airborne imagery is available.','Principle: separate independent hydric evidence from chemistry and spectra to avoid circular inference.'],.9,1.55,11.5,4.5,20)
 footer(s,'Main results: slides 1–12  |  Technical appendix: slides 13–20')
 s=slide(prs); title(s,'Study workflow keeps evidence sources independent','Each stage narrows the population without inventing a definitive hydric label')
 stages=[('KSSL audit','61,949 layers'),('Surface cohort','12,241 pedons'),('MT–ND GIS','404 points'),('MIR available','398 samples'),('Validation','Project + state holdouts')]
 for i,(a,b) in enumerate(stages):
  x=.55+i*2.55; z=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(x),Inches(2.1),Inches(2.15),Inches(1.25)); z.fill.solid(); z.fill.fore_color.rgb=PALE; z.line.color.rgb=TEAL
  text(s,x+.1,2.34,1.95,.3,a,15,NAVY,True,PP_ALIGN.CENTER); text(s,x+.1,2.8,1.95,.25,b,11,GRAY,False,PP_ALIGN.CENTER)
  if i<4: text(s,x+2.18,2.5,.3,.3,'→',20,TEAL,True,PP_ALIGN.CENTER)
 bullets(s,['Laboratory properties and MIR are evaluated against—not used to create—SSURGO/NWI evidence.','The analytical unit is a sampled layer; MIR replicates are averaged technical measurements.','Claims are tightened whenever stricter validation changes the result.'],1.0,4.15,11.3,2.1,16)
 s=slide(prs); title(s,'The KSSL archive supports a traceable analysis','Database audit established coverage, provenance, observation unit, and replicate structure')
 image(s,R/'outputs/figures/kssl_analysis/property_coverage.png',.5,1.25,7.5)
 bullets(s,['61,949 unique sampled layers','251,650 MIR scan records','Four scans are technical replicates','Measured and derived properties remain distinct','Surface layers prioritized for remote-sensing relevance'],8.35,1.5,4.15,4.8,15)
 footer(s,'Takeaway: the archive is large and usable, but method and replicate controls are essential.')
 s=slide(prs); title(s,'Laboratory properties form coherent multivariate gradients','Correlations validate extraction; PCA shows continuous structure rather than discrete hydric classes')
 image(s,R/'outputs/figures/kssl_analysis/property_correlation_heatmap.png',.4,1.25,6.35)
 image(s,R/'outputs/figures/kssl_analysis/lab_property_pca_by_depth.png',6.95,1.35,5.95)
 footer(s,'Takeaway: chemistry is interpretable, but overlapping depth classes confirm that chemistry alone is not a hydric label.')
 s=slide(prs); title(s,'Spatial overlays connect laboratory samples to mapped wetland evidence','Montana and North Dakota surface samples evaluated against independent SSURGO and NWI context')
 image(s,R/'outputs/figures/kssl_spatial_results/kssl_mt_nd_spatial_evidence_map.png',.35,1.25,9.7)
 bullets(s,['404 georeferenced surface samples','181 overlap SSURGO hydric evidence','18 intersect NWI polygons','14 are supported by both sources','Mapped evidence is a weak reference, not field-confirmed truth'],10.35,1.55,2.55,4.9,14)
 footer(s,'Takeaway: the map creates the geographic bridge needed to compare laboratory spectra with independent landscape evidence.')
 s=slide(prs); title(s,'Independent spatial evidence focuses the study on 404 surface samples','SSURGO hydric percentage and NWI intersection provide complementary regional context')
 image(s,R/'outputs/figures/kssl_spatial_results/spatial_evidence_group_counts.png',.45,1.25,6.45)
 image(s,R/'outputs/figures/kssl_spatial_results/nwi_rate_by_ssurgo_class.png',6.85,1.28,6.05)
 footer(s,'Takeaway: NWI intersections concentrate in high-SSURGO classes, but disagreements remain scientifically informative.')
 s=slide(prs); title(s,'The version-0 tool exposes evidence strength and its construction','Transparent mapped score; input completeness is reported separately from scientific certainty')
 image(s,R/'outputs/figures/kssl_tool/kssl_mt_nd_hydric_evidence_v0_diagnostic.png',.55,1.2,9.15)
 bullets(s,['374 of 404 samples fall in the provisional low-evidence category.','NWI and SSURGO remain weak reference layers.','A complete input record is not equivalent to a confirmed hydric interpretation.','The score is designed for auditability and later calibration.'],9.95,1.55,2.8,4.8,14)
 footer(s,'Takeaway: version 0 is an evidence-accounting tool, not a regulatory or field-confirmed classifier.')
 s=slide(prs); title(s,'Laboratory properties overlap across spatial-evidence groups','Exploratory comparison preserves chemistry as an independent response')
 image(s,R/'outputs/figures/kssl_spatial_results/lab_properties_by_spatial_group.png',.5,1.15,8.35)
 bullets(s,['No single property cleanly separates groups.','Spatial groups are highly unequal.','Multivariate spectroscopy is more appropriate than fixed chemistry thresholds.'],9.1,1.85,3.65,3.7,15)
 footer(s,'Takeaway: mapped wetland context cannot be reduced to one laboratory measurement.')
 s=slide(prs); title(s,'A quality-controlled MIR cohort is analysis-ready','398 samples; 1,592 selected scans; common 4000–600 cm⁻¹ grid')
 image(s,R/'outputs/figures/kssl_mir_qc/mean_mir_spectra_by_spatial_group.png',.45,1.2,8.2)
 bullets(s,['Four replicates averaged per selected master','Eleven rescans resolved without double-counting','Six missing samples documented','Replicate shapes remain highly consistent'],8.9,1.7,3.65,4.2,15)
 footer(s,'Shaded regions show within-group interquartile ranges; groups are evidence categories, not confirmed classes.')
 s=slide(prs); title(s,'MIR variation is continuous—not a simple hydric split','SNV + first-derivative PCA; unsupervised and descriptive')
 image(s,R/'outputs/figures/kssl_mir_analysis/mir_pca_by_spatial_group.png',.55,1.2,7.7)
 bullets(s,['PC1 explains 31.2%; PC2 explains 20.5%.','Spatial groups overlap strongly.','Wetland status reflects hydrology and landscape context as well as soil composition.','Supervised results require independent validation.'],8.55,1.65,4.0,4.6,15)
 s=slide(prs); title(s,'MIR strongly predicts laboratory bridge properties','Five-fold validation holds out complete KSSL projects')
 image(s,R/'outputs/figures/kssl_mir_bridge/mir_grouped_cv_summary.png',.5,1.2,7.7)
 bullets(s,['Total carbon: R² 0.92; ρ 0.91','Clay: R² 0.84; ρ 0.90','CEC: R² 0.84; ρ 0.93','pH and water retention also transfer well','Iron prediction is more modest with fewer observations'],8.5,1.45,4.15,4.9,15)
 footer(s,'Takeaway: MIR provides defensible laboratory variables for future laboratory-to-airborne comparison.')
 s=slide(prs); title(s,'Spatial hydric evidence is detectable within project-grouped validation','Ordinal weak-reference score: neither=0, one mapped source=1, both sources=2')
 image(s,R/'outputs/figures/kssl_mir_bridge/mir_grouped_cv_observed_predicted.png',.4,1.15,8.8)
 bullets(s,['Project-grouped spatial score: R² 0.17; ρ 0.53.','This is weaker than laboratory-property prediction.','The evidence score is not field-confirmed truth.','Geographic validation is required before interpretation.'],9.4,1.6,3.3,4.6,14)
 s=slide(prs); title(s,'Geographic holdouts change the hydric interpretation','Train in one state and predict only the other state')
 image(s,R/'outputs/figures/kssl_mir_geographic_validation/mir_state_holdout_spearman.png',.5,1.2,8.0)
 bullets(s,['Several laboratory-property rankings transfer across states.','Absolute calibration shifts for some properties.','The spatial evidence score does not transfer geographically.','The earlier hydric association is regional—not a general delineation model.'],8.75,1.45,3.8,5.0,15)
 s=slide(prs); title(s,'What the evidence supports—and the next bridge','Strong laboratory inference; hydric delineation still requires local surface observations')
 bullets(s,['Supported: MIR recovers key soil-property gradients under project and geographic holdouts.','Supported: SSURGO and NWI focus the cohort and expose disagreement cases.','Not supported: a universal MIR-only hydric/non-hydric classifier.','Next: relate exposed-surface airborne VNIR–SWIR spectra to MIR-derived properties.','Control vegetation, moisture, mixing, acquisition date, scale, and surface–subsurface mismatch.','Pursue field-confirmed hydric indicators for independent validation.'],.85,1.35,11.7,5.5,18)
 # Appendix
 s=slide(prs); section(s,'Technical appendix','Diagnostics, secondary figures, and interpretive safeguards')
 appendix=[
 ('A1. Depth is a first-order remote-sensing constraint','Layer midpoint classes; boxes show distributions with physical units','outputs/figures/kssl_analysis/property_distributions_by_depth.png','Airborne sensing primarily observes the exposed surface, so depth remains explicit.'),
 ('A2. Candidate evidence is not a complete hydric label','Taxonomy and supporting context before spatial overlays','outputs/figures/kssl_hydric_evidence/surface_hydric_evidence_counts.png','Indeterminate records are not automatically non-hydric.'),
 ('A3. KSSL sampling is geographically extensive','Surface points before regional filtering','outputs/figures/kssl_hydric_evidence/kssl_surface_hydric_evidence_map.png','The map shows sample distribution—not wetland boundaries.'),
 ('A4. Data coverage varies across spatial groups','Percentage of samples with each laboratory property and MIR master','outputs/figures/kssl_spatial_results/coverage_by_spatial_group.png','Missingness is retained and reported rather than silently imputed.'),
 ('A5. Technical replicate QC','Maximum replicate-to-mean RMSE for each selected master','outputs/figures/kssl_mir_qc/mir_replicate_rmse_distribution.png','High spectral correlations indicate that large RMSE cases mainly reflect intensity shifts.'),
 ('A6. MIR PCA loadings','Wavenumbers contributing to PC1 and PC2 after SNV + first derivative','outputs/figures/kssl_mir_analysis/mir_pca_loadings.png','Loadings identify influential spectral regions but are not chemical assignments by themselves.'),
 ('A7. Project-grouped observed-versus-predicted results','Every prediction comes from a fold excluding its entire KSSL project','outputs/figures/kssl_mir_bridge/mir_grouped_cv_observed_predicted.png','Strong chemistry prediction does not imply hydric-classification performance.'),
 ]
 for t,sub,p,cap in appendix:
  s=slide(prs); title(s,t,sub); image(s,R/p,.65,1.25,8.8); bullets(s,[cap],9.7,2.0,2.8,2.5,14)
 OUT.parent.mkdir(parents=True,exist_ok=True); prs.save(OUT); print(f'{OUT}\nslides={len(prs.slides)}')
if __name__=='__main__': main()

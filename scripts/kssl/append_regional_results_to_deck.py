"""Append verified Great Plains results to the current meeting-ready deck."""
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT = Path(__file__).resolve().parents[2]
SOURCE = ROOT / "outputs/presentations/KSSL_hydric_project_complete_results_MEETING_READY_v4.pptx"
OUTPUT = ROOT / "outputs/presentations/KSSL_hydric_project_complete_results_MEETING_READY_v5_REGIONAL.pptx"
NAVY = RGBColor(23, 50, 77)
TEAL = RGBColor(15, 124, 128)
GOLD = RGBColor(217, 164, 65)
GRAY = RGBColor(102, 116, 125)
PALE = RGBColor(234, 242, 243)


def textbox(slide, x, y, w, h, value, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, heading, subtitle=""):
    textbox(slide, 0.55, 0.31, 12.2, 0.55, heading, 24, NAVY, True)
    textbox(slide, 0.57, 0.90, 12.0, 0.30, subtitle, 11, GRAY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.14))
    bar.fill.solid(); bar.fill.fore_color.rgb = TEAL; bar.line.fill.background()


def bullets(slide, items, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item; p.font.name = "Aptos"; p.font.size = Pt(size); p.font.color.rgb = NAVY
        p.space_after = Pt(9)


def footer(slide, value):
    textbox(slide, 0.65, 7.08, 12.0, 0.22, value, 9, GRAY, False, PP_ALIGN.CENTER)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def main():
    prs = Presentation(SOURCE)
    # Refresh the embedded MT-ND MIR figure after the unit/title correction.
    mir_slide = prs.slides[8]
    pictures = [shape for shape in mir_slide.shapes if shape.shape_type == 13]
    if pictures:
        old_picture = max(pictures, key=lambda shape: shape.width * shape.height)
        left, top = old_picture.left, old_picture.top
        width, height = old_picture.width, old_picture.height
        old_picture._element.getparent().remove(old_picture._element)
        mir_slide.shapes.add_picture(
            str(ROOT / "outputs/figures/kssl_mir_qc/mean_mir_spectra_by_spatial_group.png"),
            left, top, width=width, height=height
        )
    slide = add_slide(prs)
    title(slide, "Regional expansion tests whether the MIR bridge generalizes", "One shallowest quality-screened layer per independent pedon in MT, ND, SD, and NE")
    stages = [("Candidate pedons", "1,321"), ("Usable MIR", "1,315"), ("KSSL projects", "214"), ("Spectral variables", "1,701")]
    for i, (label, value) in enumerate(stages):
        x = 0.75 + i * 3.1
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.55), Inches(2.55), Inches(1.25))
        shape.fill.solid(); shape.fill.fore_color.rgb = PALE; shape.line.color.rgb = TEAL
        textbox(slide, x + 0.1, 1.77, 2.35, 0.3, value, 23, TEAL, True, PP_ALIGN.CENTER)
        textbox(slide, x + 0.1, 2.28, 2.35, 0.25, label, 12, GRAY, False, PP_ALIGN.CENTER)
    bullets(slide, [
        "Technical replicate scans were quality-checked and averaged before modeling.",
        "Only six candidate pedons lacked a usable replicated MIR spectrum; two individual scans were rejected.",
        "Holding out complete projects reduces leakage from related samples collected in the same survey effort.",
        "State holdouts test a harder question: whether relationships transfer to a new geographic region."
    ], 1.0, 3.45, 11.4, 2.55, 17)
    footer(slide, "This regional cohort strengthens inference; it does not create a field-confirmed hydric label.")

    slide = add_slide(prs)
    title(slide, "Regional MIR spectra retain shared structure and state-level variation", "Median spectrum by state after averaging technical replicates; shading is the interquartile range")
    slide.shapes.add_picture(str(ROOT / "outputs/figures/kssl_regional_expansion/regional_median_mir_spectra_by_state.png"), Inches(0.55), Inches(1.25), width=Inches(8.4))
    bullets(slide, [
        "All states share the major MIR absorption structure.",
        "Differences in level and band shape reflect real variation in composition and sampling populations.",
        "The x-axis is wavenumber (cm^-1); absorbance is unitless.",
        "These curves are descriptive summaries, not hydric-soil classifications."
    ], 9.2, 1.65, 3.55, 4.6, 15)
    footer(slide, "Takeaway: the expanded library is spectrally coherent enough for regional transfer testing.")

    slide = add_slide(prs)
    title(slide, "MIR robustly recovers several laboratory bridge properties", "Five-fold validation holds out complete KSSL projects across four Great Plains states")
    slide.shapes.add_picture(str(ROOT / "outputs/figures/kssl_regional_expansion/regional_project_grouped_observed_predicted.png"), Inches(0.35), Inches(1.18), width=Inches(8.9))
    bullets(slide, [
        "Strong: total carbon, clay, pH, CEC, and 15-bar water retention.",
        "Weak: dithionite Fe, oxalate Fe, and 1/3-bar water retention.",
        "Points near the dashed 1:1 line have accurate predictions.",
        "Fixed-model results are exploratory; final models require nested tuning and spatial validation."
    ], 9.5, 1.55, 3.35, 4.9, 14)
    footer(slide, "Takeaway: prioritize the reliable MIR-derived properties as airborne comparison targets.")

    slide = add_slide(prs)
    title(slide, "Core MIR-property relationships transfer to held-out states", "Models train on three states and predict the fourth; color reports Spearman rank correlation")
    slide.shapes.add_picture(str(ROOT / "outputs/figures/kssl_regional_expansion/regional_leave_one_state_out_spearman.png"), Inches(0.55), Inches(1.20), width=Inches(8.5))
    bullets(slide, [
        "Carbon, clay, pH, CEC, and 15-bar water generally retain useful rankings in unseen states.",
        "Montana is the more difficult holdout for several targets, indicating geographic asymmetry.",
        "Extractable iron is inconsistent and should remain secondary.",
        "Next: use exposed-soil ACES pixels to test whether airborne VNIR-SWIR tracks these robust property gradients."
    ], 9.3, 1.45, 3.45, 5.1, 14)
    footer(slide, "Main claim: MIR supplies transferable soil-property information, not a universal hydric/non-hydric label.")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")
    print(f"slides={len(prs.slides)}")

if __name__ == "__main__":
    main()


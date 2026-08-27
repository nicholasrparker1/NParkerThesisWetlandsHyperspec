from pathlib import Path
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor as C
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT=Path(__file__).resolve().parents[2]; T=ROOT/'outputs/tables/kssl_analysis'; F=ROOT/'outputs/figures/kssl_analysis'
DOC=ROOT/'outputs/reports/KSSL_initial_layer_analysis.docx'; PPT=ROOT/'outputs/presentations/KSSL_initial_analysis_results.pptx'
NAVY='17324D'; TEAL='0F7C80'; GOLD='D9A441'; GRAY='66747D'; PALE='EAF2F3'

def textbox(s,x,y,w,h,text,size=12,color=NAVY,bold=False,align=PP_ALIGN.LEFT):
    b=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=b.text_frame; tf.clear(); tf.margin_left=tf.margin_right=0
    p=tf.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name='Aptos'; r.font.size=P(size); r.font.bold=bold; r.font.color.rgb=C.from_string(color); return b
def title(s,t,sub=''):
    textbox(s,.55,.35,12.2,.48,t,24,NAVY,True); textbox(s,.57,.88,12,.32,sub,11,GRAY)
    z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,I(13.333),I(.14)); z.fill.solid(); z.fill.fore_color.rgb=C.from_string(TEAL); z.line.fill.background()
def bullets(s,items,x,y,w,h,color=NAVY,size=14):
    b=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=b.text_frame; tf.clear()
    for i,item in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.level=0; p.font.name='Aptos'; p.font.size=P(size); p.font.color.rgb=C.from_string(color); p.space_after=P(11)
def add_image(s,path,x,y,w,h=None): s.shapes.add_picture(str(path),I(x),I(y),width=I(w),height=I(h) if h else None)

def word():
    desc=pd.read_csv(T/'kssl_priority_property_descriptive_statistics.csv',index_col=0); corr=pd.read_csv(T/'kssl_priority_property_top_correlations.csv'); pca=pd.read_csv(T/'kssl_lab_property_pca_variance.csv')
    d=Document(); sec=d.sections[0]; sec.top_margin=sec.bottom_margin=Inches(.65); sec.left_margin=sec.right_margin=Inches(.75)
    d.styles['Normal'].font.name='Aptos'; d.styles['Normal'].font.size=Pt(10)
    p=d.add_paragraph(); p.alignment=WD_ALIGN_PARAGRAPH.CENTER; r=p.add_run('KSSL Initial Layer-Level Analysis'); r.bold=True; r.font.size=Pt(24); r.font.color.rgb=RGBColor.from_string(NAVY)
    d.add_paragraph('Exploratory results following method-controlled database extraction').alignment=WD_ALIGN_PARAGRAPH.CENTER
    d.add_heading('Analytical population',1); d.add_paragraph('The analysis-ready table contains 61,949 unique sampled layers. Fifty-seven layers triggered non-destructive plausibility flags; descriptive and multivariate results use the remaining 61,892 layers. PCA uses 15,413 complete cases across eight priority variables. No hydric/non-hydric response label is assigned.')
    d.add_heading('Method and duplicate control',1)
    for x in ['One row represents one sampled layer.','Technical MIR scans are summarized as metadata and are not independent observations.','Each selected value retains record, definition, method, preparation, unit, and candidate-count provenance.','Where multiple candidates occurred, the globally dominant compatible method key was selected deterministically; alternatives remain in the long provenance table.']: d.add_paragraph(x,style='List Bullet')
    d.add_heading('Coverage and descriptive takeaways',1)
    table=d.add_table(rows=1,cols=4); table.style='Light Shading Accent 1';
    for c,x in zip(table.rows[0].cells,['Variable','Coverage','Median','Mean']): c.text=x
    labels={'total_carbon_pct':'Total carbon (%)','estimated_organic_carbon_pct':'Estimated organic C (%)','fe_dithionite_pct':'Dithionite Fe (%)','fe_oxalate_pct':'Oxalate Fe (%)','clay_pct':'Clay (%)','ph_water':'pH, water','water_retention_15bar_pct':'15-bar water (%)','cec_nh4oac_cmol_kg':'CEC (cmol(+)/kg)'}
    for v,l in labels.items():
        row=table.add_row().cells; row[0].text=l; row[1].text=f"{desc.loc[v,'coverage_pct']:.1f}%"; row[2].text=f"{desc.loc[v,'50%']:.2f}"; row[3].text=f"{desc.loc[v,'mean']:.2f}"
    d.add_picture(str(F/'property_coverage.png'),width=Inches(7)); d.paragraphs[-1].alignment=WD_ALIGN_PARAGRAPH.CENTER
    d.add_heading('Correlation structure',1); d.add_paragraph('The strongest associations are scientifically coherent and provide an internal validation of extraction and harmonization. They are descriptive, not causal and not independent because multiple layers may belong to one pedon.')
    table=d.add_table(rows=1,cols=3); table.style='Light Shading Accent 1';
    for c,x in zip(table.rows[0].cells,['Variable 1','Variable 2','Spearman ρ']): c.text=x
    for _,r in corr.head(12).iterrows():
        q=table.add_row().cells; q[0].text=r.variable_1; q[1].text=r.variable_2; q[2].text=f'{r.spearman_rho:.3f}'
    d.add_picture(str(F/'property_correlation_heatmap.png'),width=Inches(6.8)); d.paragraphs[-1].alignment=WD_ALIGN_PARAGRAPH.CENTER
    d.add_heading('Depth structure',1); d.add_paragraph('Depth classes describe layer midpoint, not airborne sensing depth. Carbon and related properties generally show strong depth structure, reinforcing the need to begin airborne linkage with surface or near-surface layers and to retain depth in every model.')
    d.add_picture(str(F/'property_distributions_by_depth.png'),width=Inches(7.1)); d.paragraphs[-1].alignment=WD_ALIGN_PARAGRAPH.CENTER
    d.add_heading('PCA',1); d.add_paragraph(f"PC1 explains {100*pca.iloc[0].explained_variance_ratio:.1f}% and PC2 explains {100*pca.iloc[1].explained_variance_ratio:.1f}% (cumulative {100*pca.iloc[1].cumulative_variance:.1f}%). PC1 broadly represents water-retention/CEC/carbon-nitrogen richness; PC2 contrasts texture/Fe with carbon-nitrogen. PCA is exploratory and complete-case based.")
    d.add_picture(str(F/'lab_property_pca_by_depth.png'),width=Inches(6.7)); d.paragraphs[-1].alignment=WD_ALIGN_PARAGRAPH.CENTER
    d.add_heading('Interpretive boundaries',1)
    for x in ['These results do not demonstrate hydric-soil classification performance.','KSSL samples are prepared laboratory materials, whereas airborne spectra observe mixed ambient surfaces.','Layer observations are clustered within pedons and projects; inferential modeling must use grouped validation.','Estimated organic carbon is derived and must remain distinguishable from direct measurements.','The next scientific step is defining independent hydric evidence and a surface-layer cohort.']: d.add_paragraph(x,style='List Bullet')
    d.add_heading('Reproducible products',1)
    for x in ['data/processed/kssl_layer_analysis_table.csv','outputs/tables/kssl_analysis/kssl_selected_property_provenance_long.csv','outputs/tables/kssl_analysis/','outputs/figures/kssl_analysis/','scripts/kssl/build_kssl_layer_analysis_table.py','scripts/kssl/analyze_kssl_layer_table.py']: d.add_paragraph(x)
    DOC.parent.mkdir(parents=True,exist_ok=True); d.save(DOC)

def slides():
    prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'KSSL layer-level dataset is ready for exploratory analysis','61,949 unique layers; method-controlled values with full provenance')
    cards=[('61,949','unique layers'),('61,892','quality-screened'),('19','priority properties'),('840,495','traceable values')]
    for i,(n,l) in enumerate(cards):
        x=.65+i*3.1; sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,I(x),I(1.55),I(2.75),I(1.15)); sh.fill.solid(); sh.fill.fore_color.rgb=C.from_string(PALE); sh.line.color.rgb=C.from_string('C5D8DA'); textbox(s,x+.1,1.75,2.55,.35,n,22,NAVY,True,PP_ALIGN.CENTER); textbox(s,x+.1,2.22,2.55,.25,l,10,GRAY,False,PP_ALIGN.CENTER)
    bullets(s,['One sampled layer = one analytical row','Four MIR files remain technical replicates','Multiple candidates resolved deterministically','Measured and derived properties remain distinct','No hydric label has been invented'],.85,3.25,5.6,3.3,size=16)
    add_image(s,F/'property_coverage.png',6.5,3.0,6.25,3.8)

    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Expected soil-property relationships validate the extraction','Spearman correlations among quality-screened layers')
    add_image(s,F/'property_correlation_heatmap.png',.45,1.25,6.6,5.95)
    bullets(s,['pH water ↔ pH CaCl₂: ρ = 0.98','Estimated organic C ↔ total N: ρ = 0.90','15-bar water ↔ CEC: ρ = 0.85','Clay ↔ 15-bar water: ρ = 0.79','Dithionite Fe ↔ oxalate Fe: ρ = 0.66'],7.35,1.75,5.2,3.4,size=17)
    textbox(s,7.4,5.55,4.9,.75,'Internal consistency supports moving forward—but correlations are descriptive, not causal.',13,TEAL,True)

    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Depth is a first-order constraint on airborne linkage','Boxes show quality-screened layer values by midpoint depth class')
    add_image(s,F/'property_distributions_by_depth.png',.55,1.3,8.35,5.75)
    bullets(s,['Airborne imagery primarily senses the exposed surface.','KSSL chemistry represents discrete horizons, often below the surface.','Initial linkage should prioritize layers beginning at ≤10 cm.','Depth remains a covariate—not a nuisance to discard.'],9.2,1.8,3.55,3.7,size=15)

    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Laboratory properties form interpretable multivariate gradients','PCA: 15,413 complete cases; exploratory, not a classifier')
    add_image(s,F/'lab_property_pca_by_depth.png',.6,1.3,7.4,5.7)
    bullets(s,['PC1: 43.6% of variance','PC1 loads on water retention, CEC, C/N, and Fe','PC2: 19.8% of variance','First two PCs explain 63.4% cumulatively','Depth classes overlap: chemistry alone is not a hydric label'],8.35,1.75,4.35,4.3,size=16)

    s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'What we can claim now—and what comes next','The audit supports readiness, not hydric classification performance')
    textbox(s,.75,1.45,5.5,.35,'Supported now',18,TEAL,True); bullets(s,['Large, traceable laboratory population','Strong coverage for bridge variables','Scientifically coherent property structure','Explicit replicate and method safeguards','Surface-focused cohort can be defined'],.75,1.95,5.5,4.6,size=16)
    textbox(s,6.9,1.45,5.5,.35,'Still required',18,GOLD,True); bullets(s,['Independent hydric-soil reference evidence','Pedon/project-grouped validation','Surface-layer sensitivity analysis','Selective MIR file export after cohort selection','Airborne comparison under cover/moisture controls'],6.9,1.95,5.5,4.6,size=16)
    PPT.parent.mkdir(parents=True,exist_ok=True); prs.save(PPT)

word(); slides(); print(DOC); print(PPT)

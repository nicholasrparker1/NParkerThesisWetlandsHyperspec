"""Append the validated NEON WOOD exposed-soil pilot to the main results deck."""
from pathlib import Path

import matplotlib.pyplot as plt
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[2]
SOURCE = ROOT / "outputs/presentations/KSSL_hydric_project_complete_results_MEETING_READY_v5_REGIONAL.pptx"
OUTPUT = ROOT / "outputs/presentations/KSSL_hydric_project_complete_results_MEETING_READY_v7_WOOD_VALIDATED.pptx"
RESULTS = ROOT / "outputs/tables/neon_wood_bare_soil/wood_2025_full_spectra_cross_validation.csv"
FIGURE = ROOT / "outputs/figures/neon_wood_bare_soil/wood_2025_classifier_spatial_validation.png"

NAVY = RGBColor(23, 50, 77)
TEAL = RGBColor(15, 124, 128)
GOLD = RGBColor(217, 164, 65)
GRAY = RGBColor(102, 116, 125)
PALE = RGBColor(234, 242, 243)
WHITE = RGBColor(255, 255, 255)


def textbox(slide, x, y, w, h, value, size=14, color=NAVY, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.margin_left = frame.margin_right = 0
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = value
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def title(slide, heading, subtitle=""):
    textbox(slide, 0.55, 0.31, 12.2, 0.55, heading, 24, NAVY, True)
    textbox(slide, 0.57, 0.90, 12.0, 0.30, subtitle, 11, GRAY)
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.14))
    bar.fill.solid()
    bar.fill.fore_color.rgb = TEAL
    bar.line.fill.background()


def bullets(slide, items, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = item
        p.font.name = "Aptos"
        p.font.size = Pt(size)
        p.font.color.rgb = NAVY
        p.space_after = Pt(10)


def footer(slide, value):
    textbox(slide, 0.65, 7.08, 12.0, 0.22, value, 9, GRAY, False, PP_ALIGN.CENTER)


def add_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def make_figure():
    data = pd.read_csv(RESULTS).set_index("model")
    order = ["NDVI + MNDWI", "Full spectra: PCA(10) + logistic"]
    labels = ["Two spectral indices", "Full hyperspectral signature"]
    metrics = [
        ("balanced_accuracy_mean", "Balanced accuracy"),
        ("precision_mean", "Precision"),
        ("recall_mean", "Recall"),
        ("roc_auc_mean", "ROC AUC"),
    ]
    colors = ["#D9A441", "#0F7C80"]
    plt.rcParams.update({"font.family": "Arial", "font.size": 11})
    fig, ax = plt.subplots(figsize=(10.2, 5.4))
    x = range(len(metrics))
    width = 0.34
    for i, (model, label, color) in enumerate(zip(order, labels, colors)):
        values = [float(data.loc[model, column]) for column, _ in metrics]
        positions = [v + (i - 0.5) * width for v in x]
        bars = ax.bar(positions, values, width=width, color=color, label=label)
        for bar, value in zip(bars, values):
            ax.text(bar.get_x() + bar.get_width()/2, value + 0.018, f"{value:.2f}",
                    ha="center", va="bottom", fontsize=10, color="#17324D", fontweight="bold")
    ax.set_ylim(0, 1.06)
    ax.set_ylabel("Cross-validated score")
    ax.set_xticks(list(x), [label for _, label in metrics])
    ax.grid(axis="y", alpha=0.2)
    ax.set_axisbelow(True)
    ax.spines[["top", "right"]].set_visible(False)
    ax.legend(frameon=False, loc="lower right")
    fig.tight_layout()
    FIGURE.parent.mkdir(parents=True, exist_ok=True)
    fig.savefig(FIGURE, dpi=300, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def workflow_box(slide, x, heading, detail, accent):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.58), Inches(2.12), Inches(1.48)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = PALE
    shape.line.color.rgb = accent
    textbox(slide, x + 0.12, 1.82, 1.88, 0.32, heading, 15, NAVY, True, PP_ALIGN.CENTER)
    textbox(slide, x + 0.12, 2.35, 1.88, 0.42, detail, 11, GRAY, False, PP_ALIGN.CENTER)


def main():
    prs = Presentation(SOURCE)

    slide = add_slide(prs)
    title(
        slide,
        "NEON WOOD provides a controlled rehearsal for airborne soil screening",
        "1 m surface reflectance acquired 18 June 2025; the same screening logic can be transferred to ACES",
    )
    stages = [
        ("NEON image", "426 hyperspectral bands"),
        ("Candidate mask", "NDVI + MNDWI rules"),
        ("Visual review", "150 balanced points"),
        ("Clear labels", "82 usable points"),
        ("Model test", "indices vs. full spectra"),
    ]
    accents = [TEAL, TEAL, GOLD, GOLD, TEAL]
    for i, ((heading, detail), accent) in enumerate(zip(stages, accents)):
        x = 0.48 + i * 2.55
        workflow_box(slide, x, heading, detail, accent)
        if i < len(stages) - 1:
            textbox(slide, x + 2.17, 2.05, 0.35, 0.35, "→", 23, TEAL, True, PP_ALIGN.CENTER)

    bullets(slide, [
        "Clear reference labels: 26 exposed soil, 46 vegetation, and 10 road/built points.",
        "Sixty-eight ambiguous points were retained as uncertain and excluded from model scoring.",
        "These are visual reference labels—not field observations or hydric-soil truth.",
        "This stage answers one necessary question: which airborne pixels are exposed soil?"
    ], 0.90, 3.62, 11.55, 2.62, 16)
    footer(slide, "Contribution: validates the surface-screening step required before airborne spectra can be compared with KSSL soil properties.")

    slide = add_slide(prs)
    title(
        slide,
        "Full spectra improve exposed-soil screening under spatial holdouts",
        "Five 2 km spatial-block folds; 82 clear labels from the 18 June 2025 NEON WOOD image",
    )
    slide.shapes.add_picture(str(FIGURE), Inches(0.42), Inches(1.28), width=Inches(8.15))
    bullets(slide, [
        "Balanced accuracy increases from 0.85 to 0.93; bare-soil recall is 0.92 for both models.",
        "False non-soil detections fall from 13 to 4, including road/built errors falling from 10 to 3.",
        "The full spectrum adds useful separation beyond vegetation and water indices.",
        "This remains a single-site pilot—not final transfer validation or a hydric-soil classifier."
    ], 8.82, 1.52, 4.02, 4.98, 14)
    footer(slide, "Supported claim: full airborne spectra improve exposed-soil screening at WOOD; hydric-soil classification remains untested.")

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(f"Saved {OUTPUT}")
    print(f"slides={len(prs.slides)}")
    print(f"figure={FIGURE}")


if __name__ == "__main__":
    main()






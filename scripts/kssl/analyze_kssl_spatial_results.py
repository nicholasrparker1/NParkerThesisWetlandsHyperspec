from pathlib import Path

import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from docx import Document
from docx.shared import Inches as DocInches


ROOT = Path(__file__).resolve().parents[2]
SPATIAL = ROOT / "outputs/tables/kssl_spatial/kssl_mt_nd_spatial_evidence.csv"
LAYERS = ROOT / "data/processed/kssl_layer_analysis_table.csv"
OUT_T = ROOT / "outputs/tables/kssl_spatial_results"
OUT_F = ROOT / "outputs/figures/kssl_spatial_results"
OUT_P = ROOT / "outputs/presentations/KSSL_MT_ND_spatial_results.pptx"
OUT_D = ROOT / "outputs/reports/KSSL_MT_ND_spatial_results.docx"

PROPS = {
    "total_carbon_pct": "Total carbon (%)",
    "fe_dithionite_pct": "Dithionite Fe (%)",
    "fe_oxalate_pct": "Oxalate Fe (%)",
    "clay_pct": "Clay (%)",
    "ph_water": "pH in water",
    "water_retention_15bar_pct": "15-bar water (%)",
    "cec_nh4oac_cmol_kg": "CEC (cmol(+)/kg)",
}
ORDER = ["Both NWI + SSURGO", "SSURGO only", "NWI only", "Neither"]
COLORS = {
    "Both NWI + SSURGO": "#0F7C80",
    "SSURGO only": "#D9A441",
    "NWI only": "#9B5B7B",
    "Neither": "#9AA6AD",
}


def classify(row):
    nwi = int(row["nwi_intersect"]) == 1
    ssurgo = float(row["RASTERVALU"]) > 0
    if nwi and ssurgo:
        return "Both NWI + SSURGO"
    if ssurgo:
        return "SSURGO only"
    if nwi:
        return "NWI only"
    return "Neither"


def pct(n, d):
    return 100 * n / d if d else np.nan


def build_tables(df):
    OUT_T.mkdir(parents=True, exist_ok=True)
    df.to_csv(OUT_T / "kssl_mt_nd_spatial_analysis_table.csv", index=False)

    rows = []
    for group in ORDER:
        g = df[df.spatial_evidence_group == group]
        rows.append({
            "spatial_evidence_group": group,
            "n": len(g),
            "percent": pct(len(g), len(df)),
            "montana_n": (g.state == "Montana").sum(),
            "north_dakota_n": (g.state == "North Dakota").sum(),
            "mir_available_n": (g.mir_master_count.fillna(0) > 0).sum(),
            "mir_available_pct": pct((g.mir_master_count.fillna(0) > 0).sum(), len(g)),
            "median_ssurgo_hydric_pct": g.RASTERVALU.median(),
        })
    pd.DataFrame(rows).to_csv(OUT_T / "spatial_evidence_group_summary.csv", index=False)

    ct = pd.crosstab(df.hydric_evidence_tier, df.spatial_evidence_group).reindex(columns=ORDER, fill_value=0)
    ct.to_csv(OUT_T / "kssl_tier_by_spatial_evidence.csv")

    prop_rows = []
    for var, label in PROPS.items():
        for group in ORDER:
            x = df.loc[df.spatial_evidence_group == group, var].dropna()
            prop_rows.append({
                "property": var, "label": label, "spatial_evidence_group": group,
                "n": len(x), "median": x.median(), "q25": x.quantile(.25),
                "q75": x.quantile(.75), "mean": x.mean(),
            })
    pd.DataFrame(prop_rows).to_csv(OUT_T / "laboratory_property_summary_by_spatial_group.csv", index=False)


def build_figures(df):
    OUT_F.mkdir(parents=True, exist_ok=True)
    sns.set_theme(style="whitegrid", font_scale=1.0)

    counts = df.spatial_evidence_group.value_counts().reindex(ORDER)
    fig, ax = plt.subplots(figsize=(10, 5.5))
    bars = ax.bar(ORDER, counts, color=[COLORS[x] for x in ORDER])
    ax.bar_label(bars, labels=[f"{n} ({n/len(df):.1%})" for n in counts], padding=4, fontsize=11)
    ax.set_ylabel("Surface KSSL points")
    ax.set_xlabel("")
    ax.set_title("Independent spatial evidence narrows the MT–ND KSSL cohort", weight="bold", fontsize=16)
    ax.text(0, -0.23, "SSURGO support = mapped hydric component >0%; NWI support = point intersects a mapped wetland.",
            transform=ax.transAxes, color="#596872", fontsize=10)
    ax.set_ylim(0, max(counts) * 1.18)
    sns.despine()
    fig.tight_layout()
    fig.savefig(OUT_F / "spatial_evidence_group_counts.png", dpi=220, bbox_inches="tight")
    plt.close(fig)

    tab = pd.crosstab(df.ssurgo_hydric_class, df.nwi_intersect)
    tab = tab.reindex([
        "Not Hydric", "Partially Hydric (1-25%)", "Partially Hydric (26-50%)",
        "Mostly Hydric (51-75%)", "Mostly Hydric (76-95%)", "All Hydric"
    ]).fillna(0)
    rate = 100 * tab.get(1, pd.Series(0, index=tab.index)) / tab.sum(axis=1).replace(0, np.nan)
    n = tab.sum(axis=1)
    keep = n > 0
    fig, ax = plt.subplots(figsize=(10, 5.7))
    bars = ax.bar(np.arange(keep.sum()), rate[keep], color="#0F7C80")
    ax.set_xticks(np.arange(keep.sum()), [x.replace(" Hydric", "\nHydric") for x in rate[keep].index])
    ax.set_ylabel("KSSL points intersecting NWI (%)")
    ax.set_title("NWI intersections rise in the highest SSURGO hydric classes", weight="bold", fontsize=16)
    ax.bar_label(bars, labels=[f"{int(tab.loc[i, 1] if 1 in tab else 0)}/{int(n[i])}" for i in rate[keep].index], padding=3)
    ax.set_ylim(0, max(60, np.nanmax(rate[keep]) * 1.2))
    sns.despine()
    fig.tight_layout()
    fig.savefig(OUT_F / "nwi_rate_by_ssurgo_class.png", dpi=220, bbox_inches="tight")
    plt.close(fig)

    plot_vars = ["total_carbon_pct", "clay_pct", "ph_water", "water_retention_15bar_pct"]
    fig, axes = plt.subplots(2, 2, figsize=(12, 8))
    for ax, var in zip(axes.flat, plot_vars):
        sns.boxplot(data=df, x="spatial_evidence_group", y=var, order=ORDER,
                    palette=COLORS, ax=ax, showfliers=False, hue="spatial_evidence_group", legend=False)
        ax.set_title(PROPS[var], weight="bold")
        ax.set_xlabel("")
        ax.tick_params(axis="x", rotation=18)
        ax.set_ylabel(PROPS[var])
    fig.suptitle("Laboratory properties overlap across spatial-evidence groups", weight="bold", fontsize=17)
    fig.text(.5, .01, "Exploratory distributions; unequal sample sizes and missing laboratory values. Not a classifier test.",
             ha="center", color="#596872")
    fig.tight_layout(rect=(0, .035, 1, .96))
    fig.savefig(OUT_F / "lab_properties_by_spatial_group.png", dpi=220, bbox_inches="tight")
    plt.close(fig)

    coverage = []
    for group in ORDER:
        g = df[df.spatial_evidence_group == group]
        coverage.append({
            "group": group,
            "MIR master": pct((g.mir_master_count.fillna(0) > 0).sum(), len(g)),
            **{label: pct(g[var].notna().sum(), len(g)) for var, label in PROPS.items()}
        })
    cov = pd.DataFrame(coverage).set_index("group")
    cov.to_csv(OUT_T / "data_coverage_by_spatial_group.csv")
    fig, ax = plt.subplots(figsize=(11, 5.5))
    sns.heatmap(cov, annot=True, fmt=".0f", cmap="YlGnBu", vmin=0, vmax=100,
                cbar_kws={"label": "Coverage (%)"}, ax=ax)
    ax.set_xlabel("")
    ax.set_ylabel("")
    ax.set_title("MIR linkage is strong; chemistry coverage varies by property", weight="bold", fontsize=16)
    fig.tight_layout()
    fig.savefig(OUT_F / "coverage_by_spatial_group.png", dpi=220, bbox_inches="tight")
    plt.close(fig)


def add_title(slide, title, subtitle=""):
    box = slide.shapes.add_textbox(Inches(.55), Inches(.3), Inches(12.2), Inches(.55))
    run = box.text_frame.paragraphs[0].add_run(); run.text = title; run.font.size = Pt(24); run.font.bold = True
    run.font.color.rgb = RGBColor(23, 50, 77)
    sub = slide.shapes.add_textbox(Inches(.57), Inches(.9), Inches(12), Inches(.3))
    r = sub.text_frame.paragraphs[0].add_run(); r.text = subtitle; r.font.size = Pt(11); r.font.color.rgb = RGBColor(102,116,125)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(.14))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(15,124,128); line.line.fill.background()


def add_bullets(slide, items, x, y, w, h, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.text = item
        p.font.size = Pt(size); p.font.color.rgb = RGBColor(23, 50, 77); p.space_after = Pt(10)


def build_presentation(df):
    prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Spatial context narrows the MT–ND KSSL cohort", "404 surface points linked back to Access-derived laboratory and MIR metadata")
    s.shapes.add_picture(str(OUT_F / "spatial_evidence_group_counts.png"), Inches(.55), Inches(1.35), width=Inches(7.0))
    add_bullets(s, [
        "14 points: both NWI intersection and SSURGO hydric component",
        "167: SSURGO support only; 4: NWI only",
        "219: neither mapped source supports hydric context",
        "These are evidence groups—not confirmed hydric/non-hydric labels.",
    ], 7.85, 1.65, 4.8, 4.3, 16)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Independent spatial datasets converge at high hydric percentages", "NWI intersection is uncommon overall but concentrated in high-SSURGO classes")
    s.shapes.add_picture(str(OUT_F / "nwi_rate_by_ssurgo_class.png"), Inches(.65), Inches(1.35), width=Inches(7.4))
    add_bullets(s, [
        "18 of 404 points intersect mapped NWI wetlands.",
        "8 of 16 points in the 76–95% SSURGO class intersect NWI.",
        "3 of 7 points in the All Hydric class intersect NWI.",
        "Four NWI intersections occur in SSURGO ‘Not Hydric’ units and warrant review.",
    ], 8.25, 1.7, 4.4, 4.5, 15)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "Laboratory gradients overlap across spatial-evidence groups", "Exploratory comparison; group sizes are unequal and laboratory coverage varies")
    s.shapes.add_picture(str(OUT_F / "lab_properties_by_spatial_group.png"), Inches(.5), Inches(1.2), width=Inches(8.5))
    add_bullets(s, [
        "No single laboratory property cleanly separates the groups.",
        "This supports a multivariate spectral approach rather than fixed chemistry thresholds.",
        "Inference must account for project/pedon clustering and the small NWI-supported group.",
    ], 9.25, 1.75, 3.55, 4.4, 14)

    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(s, "The cohort is ready for selective MIR analysis", "MIR files remain technical replicates; analysis unit is the sampled layer")
    s.shapes.add_picture(str(OUT_F / "coverage_by_spatial_group.png"), Inches(.55), Inches(1.35), width=Inches(7.4))
    add_bullets(s, [
        "MIR availability is audited from the Access database.",
        "Next: export spectra only for the 404-point cohort and average technical replicates.",
        "Model spatial evidence as an ordinal/weak reference—not ground truth.",
        "Use grouped validation by KSSL project or pedon.",
    ], 8.25, 1.65, 4.3, 4.7, 15)
    OUT_P.parent.mkdir(parents=True, exist_ok=True); prs.save(OUT_P)


def build_report(df):
    d = Document(); d.add_heading("KSSL Montana–North Dakota Spatial Evidence Results", 0)
    d.add_paragraph("Exploratory integration of Access-derived KSSL surface-layer records with SSURGO hydric-component percentages and National Wetlands Inventory polygons.")
    d.add_heading("Analytical cohort", 1)
    d.add_paragraph(f"The cohort contains {len(df)} surface KSSL points: {(df.state=='Montana').sum()} in Montana and {(df.state=='North Dakota').sum()} in North Dakota. Each point remains linked by lay_id to laboratory properties and MIR inventory metadata.")
    d.add_heading("Main result", 1)
    d.add_paragraph("Fourteen points receive support from both NWI and SSURGO, 167 from SSURGO only, four from NWI only, and 219 from neither source. These are spatial evidence groups, not confirmed hydric/non-hydric labels.")
    d.add_picture(str(OUT_F / "spatial_evidence_group_counts.png"), width=DocInches(6.7))
    d.add_heading("Agreement between spatial sources", 1)
    d.add_paragraph("NWI intersection rates increase sharply in the highest SSURGO hydric classes. This convergence supports using the combined evidence to prioritize samples. Disagreement cases remain scientifically important because NWI and SSURGO differ in mapping target, scale, and date.")
    d.add_picture(str(OUT_F / "nwi_rate_by_ssurgo_class.png"), width=DocInches(6.7))
    d.add_heading("Laboratory interpretation", 1)
    d.add_paragraph("Laboratory-property distributions overlap across spatial groups. No single property is treated as a hydric diagnostic, and chemistry was not used to construct the spatial evidence groups. The overlap motivates multivariate MIR analysis rather than fixed thresholds.")
    d.add_picture(str(OUT_F / "lab_properties_by_spatial_group.png"), width=DocInches(7.0))
    d.add_heading("Limitations", 1)
    for text in [
        "NWI maps wetland habitat, while SSURGO reports hydric composition of map units; neither is point-scale field truth.",
        "A point outside NWI or in a zero-percent SSURGO unit is not automatically non-hydric.",
        "Spatial and laboratory observations are clustered by pedon and project.",
        "The NWI-intersection group is small, limiting formal predictive claims.",
        "Laboratory spectra represent dried, ground samples rather than ambient airborne surfaces.",
    ]: d.add_paragraph(text, style="List Bullet")
    d.add_heading("Next step", 1)
    d.add_paragraph("Export MIR spectra for this cohort, combine technical replicates at the sample/layer level, run spectral quality control, and test whether MIR spectral structure tracks the ordinal strength of independent spatial evidence using grouped validation.")
    OUT_D.parent.mkdir(parents=True, exist_ok=True); d.save(OUT_D)


def main():
    spatial = pd.read_csv(SPATIAL)
    layers = pd.read_csv(LAYERS, usecols=["lay_id", "mir_master_count", "mir_scan_count", "mir_passed_scan_count"])
    df = spatial.merge(layers, on="lay_id", how="left", validate="one_to_one")
    df["spatial_evidence_group"] = df.apply(classify, axis=1)
    df["spatial_evidence_group"] = pd.Categorical(df.spatial_evidence_group, ORDER, ordered=True)
    build_tables(df); build_figures(df); build_presentation(df); build_report(df)
    print(df.spatial_evidence_group.value_counts().reindex(ORDER))
    print(f"MIR available: {(df.mir_master_count.fillna(0)>0).sum()}/{len(df)}")
    print(OUT_P); print(OUT_D)


if __name__ == "__main__":
    main()

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

R=Path(__file__).resolve().parents[2]
OUT=R/'outputs/presentations/KSSL_hydric_project_master_results.pptx'
NAVY=RGBColor(23,50,77); TEAL=RGBColor(15,124,128); GOLD=RGBColor(217,164,65); GRAY=RGBColor(102,116,125)
def title(s,t,sub=''):
 b=s.shapes.add_textbox(Inches(.55),Inches(.32),Inches(12.2),Inches(.55)); r=b.text_frame.paragraphs[0].add_run(); r.text=t; r.font.name='Aptos'; r.font.size=Pt(24); r.font.bold=True; r.font.color.rgb=NAVY
 b=s.shapes.add_textbox(Inches(.57),Inches(.9),Inches(12),Inches(.3)); r=b.text_frame.paragraphs[0].add_run(); r.text=sub; r.font.name='Aptos'; r.font.size=Pt(11); r.font.color.rgb=GRAY
 z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,Inches(13.333),Inches(.14)); z.fill.solid(); z.fill.fore_color.rgb=TEAL; z.line.fill.background()
def bullets(s,items,x,y,w,h,size=15):
 b=s.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=b.text_frame; tf.clear()
 for i,item in enumerate(items):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.name='Aptos'; p.font.size=Pt(size); p.font.color.rgb=NAVY; p.space_after=Pt(10)
def image(s,p,x,y,w): s.shapes.add_picture(str(p),Inches(x),Inches(y),width=Inches(w))
def main():
 prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Integrating Laboratory Soil Spectroscopy with Spatial Wetland Evidence','Montana–North Dakota KSSL analysis toward airborne hydric-soil delineation')
 bullets(s,['Goal: determine how laboratory MIR spectra and measured soil properties can improve interpretation of remotely sensed wetland soils.','Current phase: establish a defensible laboratory–spatial bridge before airborne imagery is available.','Core safeguard: mapped evidence and laboratory chemistry remain independent.'],.9,1.55,11.5,4.5,20)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'The KSSL archive supports a traceable analysis','Database audit established the observation unit, method provenance, and replicate structure')
 image(s,R/'outputs/figures/kssl_analysis/property_coverage.png',.55,1.25,7.2)
 bullets(s,['61,949 unique sampled layers','251,650 MIR scan files','Four scans are technical replicates—not independent samples','Measured and derived properties remain distinct','Surface layers are prioritized for remote-sensing relevance'],8.05,1.45,4.55,4.9,16)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Laboratory properties form coherent—but overlapping—gradients','Initial analysis validates extraction while showing chemistry alone is not a hydric label')
 image(s,R/'outputs/figures/kssl_analysis/property_correlation_heatmap.png',.45,1.2,6.7)
 image(s,R/'outputs/figures/kssl_analysis/lab_property_pca_by_depth.png',7.15,1.35,5.7)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Independent spatial evidence focuses the study on 404 surface samples','SSURGO hydric percentage and NWI intersection provide complementary context')
 image(s,R/'outputs/figures/kssl_spatial_results/spatial_evidence_group_counts.png',.55,1.25,7.3)
 bullets(s,['14 points: both NWI and SSURGO support','167: SSURGO only; 4: NWI only','219: neither mapped source','Groups are evidence tiers—not confirmed hydric/non-hydric truth'],8.15,1.65,4.4,4.4,16)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'A quality-controlled MIR cohort is ready','398 samples; 1,592 selected scans; common 4000–600 cm⁻¹ grid')
 image(s,R/'outputs/figures/kssl_mir_qc/mean_mir_spectra_by_spatial_group.png',.5,1.2,8.0)
 bullets(s,['Four technical replicates averaged per selected master','Eleven rescans resolved without double-counting','Six missing samples documented','Replicate spectral shapes are highly consistent'],8.75,1.6,3.8,4.6,15)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'MIR strongly predicts laboratory bridge properties','Five-fold validation holds out complete KSSL projects')
 image(s,R/'outputs/figures/kssl_mir_bridge/mir_grouped_cv_summary.png',.55,1.25,7.5)
 bullets(s,['Total carbon: R² 0.92; ρ 0.91','Clay: R² 0.84; ρ 0.90','CEC: R² 0.84; ρ 0.93','pH and water retention also transfer well','These properties can bridge laboratory and future airborne analyses.'],8.3,1.45,4.3,4.9,15)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Geographic holdouts reveal the boundary of the current result','Models trained in one state are tested only in the other state')
 image(s,R/'outputs/figures/kssl_mir_geographic_validation/mir_state_holdout_spearman.png',.5,1.2,8.0)
 bullets(s,['Several laboratory-property rankings transfer across states.','Absolute calibration shifts for some properties.','The spatial hydric-evidence score does not transfer geographically.','The earlier hydric association is regional—not a general delineation model.'],8.75,1.45,3.85,4.9,15)
 s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'What the evidence supports now','The laboratory bridge is strong; hydric delineation still requires local surface observations')
 bullets(s,['Supported: MIR can recover key soil-property gradients under project and geographic holdouts.','Supported: independent spatial datasets identify a focused regional cohort and disagreement cases.','Not supported: a universal MIR-only hydric/non-hydric classifier.','Next: connect exposed-surface airborne VNIR–SWIR observations to MIR-derived bridge properties.','Required controls: vegetation, moisture, mixing, acquisition date, scale, and surface–subsurface mismatch.'],.85,1.4,11.7,5.3,18)
 OUT.parent.mkdir(parents=True,exist_ok=True); prs.save(OUT); print(OUT)
if __name__=='__main__': main()

from pathlib import Path
import pandas as pd
from docx import Document
from docx.shared import Inches, Pt
from pptx import Presentation
from pptx.util import Inches as I, Pt as P
from pptx.dml.color import RGBColor as C
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

ROOT=Path(__file__).resolve().parents[2]; T=ROOT/'outputs/tables/kssl_hydric_evidence'; F=ROOT/'outputs/figures/kssl_hydric_evidence'
DOC=ROOT/'outputs/reports/KSSL_hydric_evidence_audit.docx'; PPT=ROOT/'outputs/presentations/KSSL_hydric_evidence_and_mapping.pptx'
NAVY='17324D'; TEAL='0F7C80'; GOLD='D9A441'; GRAY='66747D'; PALE='EAF2F3'

def tx(s,x,y,w,h,text,size=12,color=NAVY,bold=False,align=PP_ALIGN.LEFT):
 b=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); p=b.text_frame.paragraphs[0]; p.alignment=align; r=p.add_run(); r.text=text; r.font.name='Aptos'; r.font.size=P(size); r.font.bold=bold; r.font.color.rgb=C.from_string(color); return b
def title(s,t,sub=''):
 tx(s,.55,.35,12.2,.5,t,24,NAVY,True); tx(s,.57,.9,12,.3,sub,11,GRAY); z=s.shapes.add_shape(MSO_SHAPE.RECTANGLE,0,0,I(13.333),I(.14)); z.fill.solid(); z.fill.fore_color.rgb=C.from_string(TEAL); z.line.fill.background()
def bullets(s,items,x,y,w,h,size=15):
 b=s.shapes.add_textbox(I(x),I(y),I(w),I(h)); tf=b.text_frame; tf.clear()
 for i,item in enumerate(items):
  p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.text=item; p.font.name='Aptos'; p.font.size=P(size); p.font.color.rgb=C.from_string(NAVY); p.space_after=P(10)

layer=pd.read_csv(T/'kssl_layer_hydric_evidence.csv'); surface=pd.read_csv(T/'kssl_surface_layer_cohort.csv',low_memory=False); summary=pd.read_csv(T/'hydric_evidence_tier_summary.csv'); ss=pd.read_csv(T/'surface_hydric_evidence_tier_summary.csv'); projects=pd.read_csv(T/'project_hydric_keyword_audit.csv')

d=Document(); d.styles['Normal'].font.name='Aptos'; d.styles['Normal'].font.size=Pt(10)
d.add_heading('KSSL Hydric-Evidence Audit and Mapping Readiness',0)
d.add_paragraph('Purpose: determine what the current KSSL snapshot can support before defining hydric and non-hydric analytical groups.')
d.add_heading('Conclusion',1); d.add_paragraph('The database contains useful hydric candidates but not a universal confirmed hydric/non-hydric response. Strong taxonomy identifies 1,140 layers (218 pedons), while Aquic taxonomy, conditional subgroups, and wet-context projects provide broader supporting evidence. Most records remain indeterminate. A negative/non-hydric class cannot be assigned responsibly from absence of evidence.')
d.add_heading('Official interpretation rules',1)
for x in ['Strong taxonomic candidate: Histosols except Folists, or Histels except Folistels (NTCHS criterion 1).','Aquic taxa and selected subgroups are supporting candidates only because NTCHS criterion 2 additionally requires field indicators or other evidence.','Hydric project names provide sampling context, not a field determination.','Organic or g-suffixed horizon names are retained as context but are not used alone to assign hydric status.','Chemistry is deliberately excluded from label construction so later comparisons are not circular.']: d.add_paragraph(x,style='List Bullet')
d.add_heading('Evidence counts',1); table=d.add_table(rows=1,cols=4); table.style='Light Shading Accent 1'
for c,x in zip(table.rows[0].cells,['Tier','Layers','Pedons','Percent of layers']): c.text=x
for _,r in summary.iterrows():
 q=table.add_row().cells; q[0].text=r.hydric_evidence_tier.replace('_',' '); q[1].text=f'{r.layer_count:,}'; q[2].text=f'{r.pedon_count:,}'; q[3].text=f'{r.layer_pct:.1f}%'
d.add_heading('Surface cohort for mapping',1); d.add_paragraph(f"The uppermost sampled layer beginning at 10 cm or shallower was selected for each pedon: {len(surface):,} pedons. Of these, 8,681 have usable coordinates for ArcGIS. The surface cohort contains 211 strong taxonomic candidates, 33 candidates with multiple supporting signals, 3,081 with one supporting signal, and 8,916 indeterminate records.")
d.add_picture(str(F/'kssl_surface_hydric_evidence_map.png'),width=Inches(7));
d.add_heading('What the current data cannot establish',1)
for x in ['No verified field-indicator code is present in this snapshot.','No direct saturation, anaerobiosis, flooding-duration, or ponding-duration measurements are present.','No defensible likely-non-hydric group can be created from missing hydric evidence.','The point map shows sample distribution and candidate evidence—not wetland boundaries.']: d.add_paragraph(x,style='List Bullet')
d.add_heading('Recommended next step',1); d.add_paragraph('Load the ArcGIS-ready surface points and overlay them with the current NRCS hydric soil list/SSURGO components, National Wetlands Inventory, hydrography, terrain, and any field delineations. Use spatial overlays as additional evidence layers, not automatic truth. Select one or more regions where KSSL candidates, comparison soils, and contextual datasets overlap well.')
d.add_heading('Priority references',1)
d.add_paragraph('USDA-NRCS, Field Indicators of Hydric Soils in the United States, Version 9.3 (2026): https://www.nrcs.usda.gov/resources/guides-and-instructions/field-indicators-of-hydric-soils-in-the-united-states')
d.add_paragraph('USDA-NRCS National Technical Committee for Hydric Soils—definition, criteria, lists, and technical standard: https://www.nrcs.usda.gov/conservation-basics/soil/national-technical-committee-for-hydric-soils-ntchs')
d.add_paragraph('Berkowitz et al. (2021), Development and application of the hydric soil technical standard: https://doi.org/10.1002/saj2.20202')
d.add_paragraph('Nur et al. (2025), Mapping soil organic matter, carbon, and nitrogen in salt marshes with UAS hyperspectral imaging: https://doi.org/10.1029/2024JG008421')
d.add_paragraph('Levy and Johnson (2021), Drone-borne reflectance spectroscopy for playa soil moisture/hydroperiod: https://doi.org/10.3390/rs13051035')
d.add_heading('Products',1)
for x in ['outputs/gis/kssl_surface_hydric_evidence_points.csv','outputs/gis/kssl_surface_hydric_evidence_points.geojson','outputs/tables/kssl_hydric_evidence/kssl_layer_hydric_evidence.csv','outputs/tables/kssl_hydric_evidence/kssl_surface_layer_cohort.csv']: d.add_paragraph(x)
DOC.parent.mkdir(parents=True,exist_ok=True); d.save(DOC)

prs=Presentation(); prs.slide_width=I(13.333); prs.slide_height=I(7.5)
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Current KSSL data identify candidates—not a complete hydric label','Official NTCHS criteria guide evidence tiers; chemistry is excluded from label construction')
s.shapes.add_picture(str(F/'surface_hydric_evidence_counts.png'),I(.55),I(1.45),width=I(6.5),height=I(4.9))
bullets(s,['211 surface pedons: strong criterion-1 taxonomy','33: multiple supporting candidate signals','3,081: one supporting candidate signal','8,916: indeterminate','No defensible negative class yet'],7.35,1.65,5.2,4.5,16)
tx(s,.75,6.6,11.8,.35,'Absence of a hydric indicator in this snapshot is not evidence that a soil is non-hydric.',13,TEAL,True,PP_ALIGN.CENTER)
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'ArcGIS mapping is the correct next bridge','Surface cohort: 12,241 pedons; 8,681 currently mappable points')
s.shapes.add_picture(str(F/'kssl_surface_hydric_evidence_map.png'),I(.45),I(1.2),width=I(8.2),height=I(5.8))
bullets(s,['Map shows KSSL sampling and evidence tiers—not wetlands.','Overlay SSURGO hydric components and the National Wetlands Inventory.','Add hydrography, terrain, floodplain, and field-delineation context.','Choose regions with strong evidence and usable comparison soils.'],8.85,1.75,3.9,4.7,14)
s=prs.slides.add_slide(prs.slide_layouts[6]); title(s,'Decision: strengthen reference evidence before classification','This preserves independence between the hydric label and laboratory chemistry')
bullets(s,['Use official field indicators or hydric technical-standard observations where available.','Use taxonomy, project context, and spatial overlays as supporting evidence.','Do not use carbon, Fe, clay, pH, or water retention to create the label they will later be tested against.','After regional selection, compare laboratory properties and selectively export raw MIR spectra.'],.9,1.55,11.7,4.9,18)
PPT.parent.mkdir(parents=True,exist_ok=True); prs.save(PPT)
print(DOC); print(PPT)

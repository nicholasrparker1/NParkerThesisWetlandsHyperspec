from __future__ import annotations

import csv
from pathlib import Path

import matplotlib.pyplot as plt
from docx import Document
from docx.enum.section import WD_ORIENT
from docx.enum.table import WD_CELL_VERTICAL_ALIGNMENT
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import Inches, Pt, RGBColor
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches as PInches, Pt as PPt


ROOT = Path(__file__).resolve().parents[2]
AUDIT = ROOT / "outputs" / "tables" / "kssl_audit"
DOC_OUT = ROOT / "outputs" / "reports" / "KSSL_database_audit.docx"
SLIDE_OUT = ROOT / "outputs" / "presentations" / "KSSL_audit_data_readiness_slide.pptx"
FIG_OUT = ROOT / "outputs" / "figures" / "kssl_priority_property_coverage.png"

NAVY = "17324D"
TEAL = "0F7C80"
GOLD = "D9A441"
PALE = "EAF2F3"
INK = "24313A"
GRAY = "66747D"


def rows(name: str) -> list[dict[str, str]]:
    with (AUDIT / name).open(encoding="utf-8-sig", newline="") as f:
        return list(csv.DictReader(f))


def integer(value: str | None) -> int:
    return int(float(value or 0))


def shade(cell, fill: str) -> None:
    tc_pr = cell._tc.get_or_add_tcPr()
    shd = tc_pr.find(qn("w:shd"))
    if shd is None:
        shd = OxmlElement("w:shd")
        tc_pr.append(shd)
    shd.set(qn("w:fill"), fill)


def set_cell_text(cell, text: str, *, bold=False, color=INK, size=9) -> None:
    cell.text = ""
    p = cell.paragraphs[0]
    run = p.add_run(str(text))
    run.bold = bold
    run.font.name = "Aptos"
    run.font.size = Pt(size)
    run.font.color.rgb = RGBColor.from_string(color)
    cell.vertical_alignment = WD_CELL_VERTICAL_ALIGNMENT.CENTER


def add_table(doc: Document, headers: list[str], data: list[list[str]], widths=None):
    table = doc.add_table(rows=1, cols=len(headers))
    table.style = "Table Grid"
    table.autofit = False
    for i, header in enumerate(headers):
        shade(table.rows[0].cells[i], NAVY)
        set_cell_text(table.rows[0].cells[i], header, bold=True, color="FFFFFF", size=9)
    for ridx, values in enumerate(data):
        cells = table.add_row().cells
        for i, value in enumerate(values):
            if ridx % 2:
                shade(cells[i], "F4F7F8")
            set_cell_text(cells[i], value)
    if widths:
        for row in table.rows:
            for i, width in enumerate(widths):
                row.cells[i].width = Inches(width)
    doc.add_paragraph()
    return table


def add_heading(doc: Document, text: str, level: int = 1) -> None:
    p = doc.add_heading(text, level=level)
    p.paragraph_format.space_before = Pt(10)
    p.paragraph_format.space_after = Pt(5)
    for run in p.runs:
        run.font.name = "Aptos Display"
        run.font.color.rgb = RGBColor.from_string(NAVY if level == 1 else TEAL)


def add_bullets(doc: Document, items: list[str]) -> None:
    for item in items:
        p = doc.add_paragraph(style="List Bullet")
        p.paragraph_format.space_after = Pt(3)
        p.add_run(item)


def make_chart(coverage: list[dict[str, str]]) -> None:
    wanted = [
        "Carbon, Total",
        "Estimated Organic Carbon",
        "Clay",
        "pH, 1:1 Soil-Water Suspension",
        "Water Retention, 15 Bar, <2mm,  Air-dry",
        "Iron, Dithionite Citrate Extractable",
        "Iron, Oxalate Extractable",
        "Bulk Density, <2mm Fraction, Ovendry",
    ]
    labels = {
        "Carbon, Total": "Total carbon",
        "Estimated Organic Carbon": "Estimated organic carbon",
        "Clay": "Clay",
        "pH, 1:1 Soil-Water Suspension": "pH (water)",
        "Water Retention, 15 Bar, <2mm,  Air-dry": "15-bar water retention",
        "Iron, Dithionite Citrate Extractable": "Dithionite-extractable Fe",
        "Iron, Oxalate Extractable": "Oxalate-extractable Fe",
        "Bulk Density, <2mm Fraction, Ovendry": "Oven-dry bulk density",
    }
    lookup = {r["property_name"]: integer(r["unique_layers"]) for r in coverage}
    values = [lookup[x] for x in wanted][::-1]
    names = [labels[x] for x in wanted][::-1]
    FIG_OUT.parent.mkdir(parents=True, exist_ok=True)
    fig, ax = plt.subplots(figsize=(9.2, 4.8))
    bars = ax.barh(names, values, color="#0F7C80")
    ax.set_xlim(0, 66000)
    ax.set_xlabel("Unique sampled layers", color="#24313A")
    ax.spines[["top", "right", "left"]].set_visible(False)
    ax.grid(axis="x", color="#D9E2E6", linewidth=0.8)
    ax.set_axisbelow(True)
    ax.tick_params(axis="y", length=0, labelcolor="#24313A")
    ax.tick_params(axis="x", colors="#66747D")
    for bar, value in zip(bars, values):
        ax.text(value + 900, bar.get_y() + bar.get_height() / 2, f"{value:,}", va="center", fontsize=9, color="#24313A")
    fig.tight_layout()
    fig.savefig(FIG_OUT, dpi=220, bbox_inches="tight", facecolor="white")
    plt.close(fig)


def build_word() -> None:
    counts = {r["table_name"]: integer(r["row_count"]) for r in rows("table_counts.csv")}
    coverage = rows("property_coverage.csv")
    context = rows("sample_context_coverage.csv")[0]
    masters = {integer(r["masters_per_sample"]): integer(r["sample_count"]) for r in rows("mir_masters_per_sample.csv")}
    scans = {integer(r["scans_per_master"]): integer(r["master_count"]) for r in rows("mir_scans_per_master.csv")}
    statuses = rows("mir_scan_status.csv")
    method_rows = rows("property_method_breakdown.csv")

    DOC_OUT.parent.mkdir(parents=True, exist_ok=True)
    doc = Document()
    section = doc.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    section.left_margin = Inches(0.7)
    section.right_margin = Inches(0.7)

    styles = doc.styles
    styles["Normal"].font.name = "Aptos"
    styles["Normal"].font.size = Pt(10)
    styles["Normal"].font.color.rgb = RGBColor.from_string(INK)

    title = doc.add_paragraph()
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run("KSSL MIR Database Audit")
    run.bold = True
    run.font.name = "Aptos Display"
    run.font.size = Pt(25)
    run.font.color.rgb = RGBColor.from_string(NAVY)
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    sr = subtitle.add_run("Data readiness for hydric-soil spectroscopy analysis")
    sr.font.size = Pt(13)
    sr.font.color.rgb = RGBColor.from_string(TEAL)
    doc.add_paragraph("Source: KSSL MIR Spectra Access Portable v2, 3 July 2025 snapshot")

    add_heading(doc, "Executive conclusion")
    p = doc.add_paragraph()
    p.add_run("The audit is slide-worthy because it establishes analytical credibility. ").bold = True
    p.add_run(
        "The database provides strong coverage for carbon, texture, pH, water retention, bulk density, and iron forms, "
        "but its scan-level structure can inflate sample size unless technical replicates, rescans, geographic links, "
        "and historical methods are handled explicitly."
    )
    add_bullets(doc, [
        "Use the sampled layer—not the individual scan filename—as the initial independent observational unit.",
        "Treat the usual four scans as technical replicates within an MIR master record.",
        "Apply an explicit policy to 955 samples with multiple MIR master records.",
        "Do not silently pool the 285 observed method/preparation/instrument/reliability combinations.",
        "Proceed to method-controlled layer-level analysis; the 100 GB raw spectral library is not yet required.",
    ])

    add_heading(doc, "1. Database scale and hierarchy")
    add_table(doc, ["Entity", "Records", "Interpretation"], [
        ["Sites", f"{counts['lims_site']:,}", "Sampling locations"],
        ["Pedons", f"{counts['lims_pedon']:,}", "Profile-level grouping for leakage-safe validation"],
        ["Layers", f"{counts['layer']:,}", "Primary initial analytical unit"],
        ["Samples", f"{counts['sample']:,}", "Physical samples associated with layers"],
        ["MIR master records", f"{counts['mir_scan_mas_data']:,}", "Scan events/batches linked to samples"],
        ["MIR scan-detail records", f"{counts['mir_scan_det_data']:,}", "Technical replicate files"],
        ["Measured analyte rows", f"{counts['layer_analyte']:,}", "Laboratory measurements"],
        ["Calculated-result rows", f"{counts['result']:,}", "Derived and predicted results"],
    ], [1.6, 1.1, 4.5])
    doc.add_paragraph(
        "The near equality of layers and samples reflects a largely one-sample-per-layer design: only six layers have two samples. "
        "Every sample occurs in the MIR master table, confirming that this portable database is a MIR-selected snapshot rather than the full national characterization population."
    )

    add_heading(doc, "2. MIR replicate, rescan, and QC structure")
    passed = sum(integer(r["scan_count"]) for r in statuses if r["qc_file_status"] == "Passed")
    blank = sum(integer(r["scan_count"]) for r in statuses if not r["qc_file_status"])
    add_table(doc, ["Audit finding", "Count", "Consequence"], [
        ["Masters with four scans", f"{scans.get(4, 0):,}", "Expected technical-replicate structure"],
        ["Masters with one scan", f"{scans.get(1, 0):,}", "Flag for spectrum-level QC"],
        ["Samples with one master", f"{masters.get(1, 0):,}", "Straightforward replicate averaging"],
        ["Samples with multiple masters", f"{masters.get(2, 0) + masters.get(3, 0):,}", "Requires rescan selection/combination rule"],
        ["Validated scans with Passed QC", f"{passed:,}", "Highest-confidence scan subset"],
        ["Validated scans with blank QC", f"{blank:,}", "Do not assume failure; preserve and investigate"],
    ], [2.25, 1.2, 3.8])
    doc.add_paragraph(
        "A canned query that joins one laboratory value to scan-detail records will normally return that value four times. "
        "Those rows are repeated file associations, not four independent measurements of soil chemistry."
    )

    add_heading(doc, "3. Context completeness")
    sample_n = integer(context["sample_rows"])
    context_data = []
    for label, key in [
        ("Top depth", "with_top_depth"),
        ("Bottom depth", "with_bottom_depth"),
        ("Horizon designation", "with_horizon"),
        ("Standardized coordinates", "with_standard_coordinates"),
        ("User pedon identifier", "with_user_pedon_id"),
    ]:
        n = integer(context[key])
        context_data.append([label, f"{n:,}", f"{100*n/sample_n:.1f}%", "Retain missingness explicitly"])
    add_table(doc, ["Context field", "Samples", "Coverage", "Treatment"], context_data, [2.2, 1.1, 1.0, 3.0])

    add_heading(doc, "4. Hydric-relevant property readiness")
    make_chart(coverage)
    doc.add_picture(str(FIG_OUT), width=Inches(7.2))
    doc.paragraphs[-1].alignment = WD_ALIGN_PARAGRAPH.CENTER
    doc.add_paragraph(
        "Coverage is counted by unique sampled layer after applying the IDs and preparation filters in the database's 36-property reference table. "
        "These counts support multivariable exploratory analysis and method-controlled modeling."
    )

    priority_names = [
        "Carbon, Total", "Estimated Organic Carbon", "Clay",
        "pH, 1:1 Soil-Water Suspension", "Water Retention, 15 Bar, <2mm,  Air-dry",
        "Iron, Dithionite Citrate Extractable", "Iron, Oxalate Extractable",
        "Bulk Density, <2mm Fraction, Ovendry", "Carbonate, <2mm Fraction",
    ]
    lookup = {r["property_name"]: r for r in coverage}
    role = {
        "Carbon, Total": "Organic matter / carbon accumulation",
        "Estimated Organic Carbon": "Organic component; derived value",
        "Clay": "Texture, water retention, and spectral mixing",
        "pH, 1:1 Soil-Water Suspension": "Soil chemical environment",
        "Water Retention, 15 Bar, <2mm,  Air-dry": "Water-retention behavior",
        "Iron, Dithionite Citrate Extractable": "Pedogenic free Fe pool",
        "Iron, Oxalate Extractable": "Poorly crystalline Fe pool",
        "Bulk Density, <2mm Fraction, Ovendry": "Organic/mineral structure and porosity",
        "Carbonate, <2mm Fraction": "Mineralogy and pH stratification",
    }
    add_table(doc, ["Priority property", "Source", "Unique layers", "Analytical role"], [
        [name, lookup[name]["route"].title(), f"{integer(lookup[name]['unique_layers']):,}", role[name]]
        for name in priority_names
    ], [2.45, 0.8, 1.0, 3.0])

    add_heading(doc, "5. Method heterogeneity and harmonization")
    doc.add_paragraph(
        f"The 35 populated reference properties resolve to {len(method_rows):,} combinations of analytical definition, preparation, "
        "size fraction, instrument set, laboratory, and reliability. This is provenance—not nuisance metadata."
    )
    add_bullets(doc, [
        "Select a dominant, scientifically compatible method family for each initial response variable.",
        "Keep method and preparation identifiers in the long-form source table.",
        "Do not combine units or extraction procedures solely because analyte names are similar.",
        "Use SSIR 42 Parts 1 and 2 to distinguish current and obsolete methods; use SSIR 45 for reporting and derived-value interpretation.",
        "KCl-extractable Fe (analyte 730) has no result rows and must be excluded unless another release supplies it.",
    ])

    add_heading(doc, "6. Biases and limitations established by the audit")
    add_table(doc, ["Issue", "Risk", "Required safeguard"], [
        ["Four scan files per master", "Fourfold pseudo-replication", "Aggregate within MIR master after QC"],
        ["Multiple masters for 955 samples", "Rescan-dependent spectra", "Document selection or averaging policy"],
        ["Multiple area links per site", "Row multiplication", "Filter each join by area type"],
        ["Taxonomy history", "Several classifications per pedon", "Use latest or declared classification rule"],
        ["Method heterogeneity", "Incompatible values pooled", "Method-controlled extraction/harmonization"],
        ["Missing hydric label", "Circular or proxy-based validation", "Obtain independent hydric evidence"],
        ["Prepared laboratory soil", "Mismatch with airborne surface", "Frame properties as bridge variables"],
        ["Subsurface layers", "Airborne sensor cannot observe depth directly", "Start with surface/near-surface subset"],
    ], [1.85, 2.15, 3.3])

    add_heading(doc, "7. Decisions supported by this audit")
    add_bullets(doc, [
        "Proceed to a method-controlled, one-row-per-layer analysis table.",
        "Prioritize carbon, texture, Fe pools, pH, water retention, and bulk density.",
        "Group cross-validation by pedon, and later by project or geography where appropriate.",
        "Separate measured, derived, MIR-predicted, and airborne variables in every result and figure.",
        "Do not retrieve the complete raw MIR library yet; query-selected exports can follow sample selection.",
    ])

    add_heading(doc, "8. Why this merits one slide")
    doc.add_paragraph(
        "The audit is not a scientific result about hydric soils. It is a data-readiness result showing that the study has enough relevant laboratory observations, "
        "and that its analysis design explicitly prevents pseudo-replication and method mixing. One slide is appropriate in an early research update or methods presentation; "
        "the full audit belongs in project documentation or supplementary methods."
    )

    doc.add_page_break()
    add_heading(doc, "Appendix A. Complete curated-property coverage")
    appendix = sorted(coverage, key=lambda r: integer(r["unique_layers"]), reverse=True)
    add_table(doc, ["Property", "Route", "Result rows", "Unique layers", "Layers with MIR"], [
        [r["property_name"], r["route"].title(), f"{integer(r['result_rows']):,}",
         f"{integer(r['unique_layers']):,}", f"{integer(r['layers_with_mir']):,}"]
        for r in appendix
    ], [3.4, 0.8, 1.0, 1.0, 1.0])

    add_heading(doc, "Appendix B. Reproducibility")
    doc.add_paragraph("Audit script: scripts/kssl/audit_kssl_database.ps1")
    doc.add_paragraph("Audit tables: outputs/tables/kssl_audit/")
    doc.add_paragraph("Database: data/raw/KSSL/MIR Spectra_Access_Portable.accdb")
    doc.add_paragraph("Fixed method references: data/raw/KSSL/documentation/")

    doc.save(DOC_OUT)


def add_textbox(slide, x, y, w, h, text, size, color, *, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(PInches(x), PInches(y), PInches(w), PInches(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = tf.margin_right = 0
    tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = "Aptos"
    r.font.size = PPt(size)
    r.font.bold = bold
    r.font.color.rgb = PptRGB.from_string(color)
    return box


def card(slide, x, y, w, h, number, label):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(x), PInches(y), PInches(w), PInches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = PptRGB.from_string(PALE)
    shape.line.color.rgb = PptRGB.from_string("C5D8DA")
    add_textbox(slide, x + .12, y + .12, w - .24, .34, number, 20, NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + .12, y + .52, w - .24, .35, label, 9, GRAY, align=PP_ALIGN.CENTER)


def build_slide() -> None:
    prs = Presentation()
    prs.slide_width = PInches(13.333)
    prs.slide_height = PInches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = PptRGB(255, 255, 255)

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, PInches(.16))
    band.fill.solid(); band.fill.fore_color.rgb = PptRGB.from_string(TEAL); band.line.fill.background()
    add_textbox(slide, .55, .38, 12.2, .48, "KSSL audit confirms a strong—but structured—analysis foundation", 25, NAVY, bold=True)
    add_textbox(slide, .57, .92, 12.0, .30, "The usable unit is the sampled layer; scan files are technical replicates, not independent soils.", 12, GRAY)

    card(slide, .60, 1.40, 1.85, 1.02, "61,955", "samples")
    card(slide, 2.58, 1.40, 1.85, 1.02, "251,650", "validated scans")
    card(slide, 4.56, 1.40, 1.85, 1.02, "35 / 36", "properties populated")
    card(slide, 6.54, 1.40, 1.85, 1.02, "955", "samples rescanned")

    add_textbox(slide, .62, 2.75, 7.8, .3, "Priority bridge variables have substantial unique-layer coverage", 14, NAVY, bold=True)
    slide.shapes.add_picture(str(FIG_OUT), PInches(.48), PInches(3.08), width=PInches(8.05), height=PInches(3.72))

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, PInches(8.72), PInches(1.40), PInches(4.02), PInches(5.40))
    panel.fill.solid(); panel.fill.fore_color.rgb = PptRGB.from_string(NAVY)
    panel.line.fill.background()
    add_textbox(slide, 9.05, 1.75, 3.35, .4, "What the audit changes", 17, "FFFFFF", bold=True)
    bullets = [
        ("1", "Aggregate the usual four scans within each MIR master."),
        ("2", "Resolve multiple masters before spectrum-level modeling."),
        ("3", "Control methods, preparations, units, and reliability."),
        ("4", "Model at layer level and group validation by pedon."),
        ("5", "Treat chemistry as a bridge—not direct airborne truth."),
    ]
    y = 2.38
    for number, text in bullets:
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, PInches(9.05), PInches(y), PInches(.34), PInches(.34))
        circ.fill.solid(); circ.fill.fore_color.rgb = PptRGB.from_string(GOLD); circ.line.fill.background()
        tf = circ.text_frame; tf.clear(); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = number; r.font.bold = True; r.font.size = PPt(10); r.font.color.rgb = PptRGB.from_string(NAVY)
        add_textbox(slide, 9.55, y - .01, 2.75, .55, text, 11, "FFFFFF")
        y += .82
    add_textbox(slide, 9.05, 6.42, 3.25, .25, "Next: method-controlled one-row-per-layer table", 10, "A9DADB", bold=True)

    SLIDE_OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(SLIDE_OUT)


if __name__ == "__main__":
    coverage = rows("property_coverage.csv")
    make_chart(coverage)
    build_word()
    build_slide()
    print(DOC_OUT)
    print(SLIDE_OUT)
    print(FIG_OUT)
